"""Extract BeraterprofilContent from an existing ORBIT Beraterprofil PPTX."""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation

from src.generator.pptx_layout_profile import LAYOUT_PROFILE, TEMPLATE_ANCHORS
from src.generator.pptx_shapes import find_shape_at
from src.models.schemas import BeraterprofilContent, CategorizedBullet, ToolCategory

_CONTENT_KEYS = (
    "title",
    "position",
    "schwerpunkte",
    "summary",
    "kompetenzen",
    "relevante",
    "international",
    "education",
    "tools",
)


def parse_beraterprofil_pptx(path: str | Path) -> BeraterprofilContent:
    """Read profile fields from a Beraterprofil slide (template or generated)."""
    prs = Presentation(str(path))
    if not prs.slides:
        raise ValueError("PPTX enthält keine Folien")

    slide = prs.slides[0]
    shapes = _locate_content_shapes(slide)

    title = _plain_text(shapes["title"])
    domain = _domain_from_title(title)
    kompetenzen = _bullet_lines(shapes["kompetenzen"])
    relevante = _parse_categorized(shapes["relevante"])
    international = _bullet_lines(shapes["international"])
    education = _bullet_lines(shapes["education"])
    tools = _parse_tool_categories(shapes["tools"])

    return BeraterprofilContent(
        domain=domain,
        title=title or f"Beraterprofil – {domain}",
        position_level=_plain_text(shapes["position"]),
        schwerpunkte=_plain_text(shapes["schwerpunkte"]),
        summary=_plain_text(shapes["summary"]),
        kompetenzen=kompetenzen,
        relevante_erfahrungen=relevante,
        international_experience=international,
        tool_categories=tools,
        education_certificates=education,
        audit_warnings=["Importiert aus hochgeladener PPTX"],
    )


def _locate_content_shapes(slide) -> dict:
    shapes: dict = {}
    for key in _CONTENT_KEYS:
        if key in LAYOUT_PROFILE:
            rect = LAYOUT_PROFILE[key]
            try:
                shapes[key] = find_shape_at(slide, rect.left, rect.top)
                continue
            except KeyError:
                pass
        left, top = TEMPLATE_ANCHORS[key]
        shapes[key] = find_shape_at(slide, left, top)
    return shapes


def _plain_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def _bullet_lines(shape) -> list[str]:
    if not shape.has_text_frame:
        return []
    lines: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            lines.append(text)
    return lines


def _parse_categorized(shape) -> list[CategorizedBullet]:
    bullets: list[CategorizedBullet] = []
    for line in _bullet_lines(shape):
        if ": " in line:
            category, details = line.split(": ", 1)
            bullets.append(CategorizedBullet(category=category.strip(), details=details.strip()))
        else:
            bullets.append(CategorizedBullet(category=line, details=""))
    return bullets


def _parse_tool_categories(shape) -> list[ToolCategory]:
    categories: list[ToolCategory] = []
    for line in _bullet_lines(shape):
        if ": " in line:
            category, tools_str = line.split(": ", 1)
            tools = [tool.strip() for tool in re.split(r",|;", tools_str) if tool.strip()]
            categories.append(ToolCategory(category=category.strip(), tools=tools))
    return categories


def _domain_from_title(title: str) -> str:
    for sep in ("–", "-", "—"):
        if sep in title:
            return title.split(sep, 1)[1].strip() or "IT-Beratung"
    return "IT-Beratung"
