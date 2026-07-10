"""Low-level PowerPoint text helpers that preserve template paragraph formatting."""

from __future__ import annotations

from copy import deepcopy

from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement


def _tx_body(text_frame):
    return text_frame._txBody


def _run_properties(run_element) -> OxmlElement | None:
    r_pr = run_element.find(qn("a:rPr"))
    return deepcopy(r_pr) if r_pr is not None else None


def _make_run(text: str, r_pr: OxmlElement | None = None) -> OxmlElement:
    run = OxmlElement("a:r")
    if r_pr is not None:
        run.append(deepcopy(r_pr))
    text_element = OxmlElement("a:t")
    text_element.text = text
    run.append(text_element)
    return run


def _replace_runs(paragraph, runs: list[OxmlElement]) -> None:
    for run in list(paragraph.findall(qn("a:r"))):
        paragraph.remove(run)
    for run in runs:
        paragraph.append(run)


def _trim_paragraphs(text_frame, keep: int) -> None:
    body = _tx_body(text_frame)
    paragraphs = body.findall(qn("a:p"))
    for paragraph in paragraphs[keep:]:
        body.remove(paragraph)


def _enable_norm_autofit(text_frame) -> None:
    """Allow text to shrink inside small boxes (e.g. Schwerpunkte)."""
    body_pr = _tx_body(text_frame).find(qn("a:bodyPr"))
    if body_pr is None:
        return
    for child in list(body_pr):
        if child.tag in {qn("a:spAutoFit"), qn("a:normAutofit"), qn("a:noAutofit")}:
            body_pr.remove(child)
    body_pr.append(OxmlElement("a:normAutofit"))


def set_minimal_plain_text(text_frame, text: str) -> None:
    """Update only run text — preserve all template paragraph/run formatting."""
    if not text_frame.paragraphs:
        return
    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = text
    _trim_paragraphs(text_frame, 1)


def set_minimal_bullets(text_frame, lines: list[str]) -> None:
    """Update bullet lines in-place without cloning or replacing XML structure."""
    if not lines:
        return

    paragraphs = text_frame.paragraphs
    for index, line in enumerate(lines):
        if index >= len(paragraphs):
            break
        paragraph = paragraphs[index]
        if paragraph.runs:
            paragraph.runs[0].text = line
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = line

    for index in range(len(lines), len(paragraphs)):
        paragraph = paragraphs[index]
        if paragraph.runs:
            paragraph.runs[0].text = ""
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = ""


def set_plain_text(text_frame, text: str, prototype) -> None:
    """Single paragraph, no bullet — used for summary and position lines."""
    proto_runs = prototype.findall(qn("a:r"))
    r_pr = _run_properties(proto_runs[0]) if proto_runs else None

    paragraphs = text_frame.paragraphs
    if paragraphs:
        _replace_runs(paragraphs[0]._p, [_make_run(text, r_pr)])
        _trim_paragraphs(text_frame, 1)
    else:
        paragraph = deepcopy(prototype)
        _replace_runs(paragraph, [_make_run(text, r_pr)])
        _tx_body(text_frame).append(paragraph)


def set_bullet_lines(text_frame, lines: list[str], prototype) -> None:
    """Update bullets in-place so each line keeps its own template paragraph properties."""
    if not lines:
        return

    proto_runs = prototype.findall(qn("a:r"))
    r_pr = _run_properties(proto_runs[0]) if proto_runs else None
    existing = list(text_frame.paragraphs)

    for index, line in enumerate(lines):
        if index < len(existing):
            paragraph = existing[index]._p
        else:
            paragraph = deepcopy(prototype)
            _tx_body(text_frame).append(paragraph)
        _replace_runs(paragraph, [_make_run(line, r_pr)])

    _trim_paragraphs(text_frame, len(lines))


def set_categorized_lines(
    text_frame,
    items: list[tuple[str, str]],
    prototype,
) -> None:
    """Category + details lines with bold/normal runs."""
    if not items:
        return

    proto_runs = prototype.findall(qn("a:r"))
    bold_r_pr = _run_properties(proto_runs[0]) if proto_runs else None
    normal_r_pr = _run_properties(proto_runs[-1]) if proto_runs else bold_r_pr
    existing = list(text_frame.paragraphs)

    for index, (category, details) in enumerate(items):
        if index < len(existing):
            paragraph = existing[index]._p
        else:
            paragraph = deepcopy(prototype)
            _tx_body(text_frame).append(paragraph)
        _replace_runs(
            paragraph,
            [
                _make_run(category, bold_r_pr),
                _make_run(f": {details}", normal_r_pr),
            ],
        )

    _trim_paragraphs(text_frame, len(items))
