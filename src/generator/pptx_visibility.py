"""Post-processing helpers to ensure text is visible in the ORBIT template."""

from __future__ import annotations

from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt


def move_shape_to_front(shape) -> None:
    """Raise shape z-order so text is not hidden behind other slide elements."""
    element = shape._element
    parent = element.getparent()
    parent.remove(element)
    parent.append(element)


def ensure_theme_text_color(text_frame) -> None:
    """Rectangle placeholders may lose font color — force readable theme text."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
            if run.font.size is None:
                run.font.size = Pt(10)


def expand_shape_height(shape, min_height_emu: int) -> None:
    if shape.height < min_height_emu:
        shape.height = min_height_emu


def enable_word_wrap(text_frame) -> None:
    text_frame.word_wrap = True
