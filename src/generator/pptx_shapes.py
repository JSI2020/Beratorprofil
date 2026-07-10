"""Locate ORBIT template shapes and apply the fine-tuned layout profile."""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.generator.pptx_layout_profile import (
    LAYOUT_PROFILE,
    TEMPLATE_ANCHORS,
    _TOLERANCE,
    ShapeRect,
)


def find_shape_at(slide, left: int, top: int, tolerance: int = _TOLERANCE):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if abs(shape.left - left) < tolerance and abs(shape.top - top) < tolerance:
            return shape
    raise KeyError(f"No text shape at ({left}, {top})")


def find_all_content_shapes(slide) -> dict:
    """Locate shapes using original template anchors."""
    return {key: find_shape_at(slide, left, top) for key, (left, top) in TEMPLATE_ANCHORS.items()}


def apply_rect(shape, rect: ShapeRect) -> None:
    shape.left = rect.left
    shape.top = rect.top
    shape.width = rect.width
    shape.height = rect.height


def apply_layout_profile(shapes: dict) -> None:
    """Move/resize text boxes to match the user-tuned layout."""
    for key, rect in LAYOUT_PROFILE.items():
        if key in shapes:
            apply_rect(shapes[key], rect)


def find_shape(slide, key: str):
    """Find a shape after layout has been applied (uses tuned coordinates)."""
    if key not in LAYOUT_PROFILE:
        raise KeyError(f"Unknown shape key: {key}")
    rect = LAYOUT_PROFILE[key]
    return find_shape_at(slide, rect.left, rect.top)


def find_right_column_divider(slide):
    from src.generator.pptx_layout_profile import RIGHT_COLUMN_LINE_LEFT, RIGHT_COLUMN_LINE_TOP_HINT

    best = None
    best_dist = None
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.LINE:
            continue
        if shape.left < 7_000_000:
            continue
        dist = abs(shape.top - RIGHT_COLUMN_LINE_TOP_HINT) + abs(shape.left - RIGHT_COLUMN_LINE_LEFT)
        if best_dist is None or dist < best_dist:
            best = shape
            best_dist = dist
    return best
