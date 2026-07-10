"""Fill the ORBIT Beraterprofil PowerPoint template."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from src.models.schemas import BeraterprofilContent

# Shape indices on slide 1 of Beraterprofil_TEMPLATE.pptx
SHAPE_TITLE = 1
SHAPE_KOMPETENZEN = 2
SHAPE_RELEVANTE = 3
SHAPE_INTERNATIONAL = 4
SHAPE_POSITION_LINE1 = 12
SHAPE_POSITION_LINE2 = 13
SHAPE_SUMMARY = 14
SHAPE_EDUCATION = 15
SHAPE_PHOTO = 17
SHAPE_TOOLS = 20

MIN_FONT_PT = 9
BODY_FONT_PT = 10


def generate_pptx(
    content: BeraterprofilContent,
    template_path: str | Path,
    output_path: str | Path,
    photo_path: str | Path | None = None,
) -> Path:
    template_path = Path(template_path)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(template_path))
    slide = prs.slides[0]

    _set_text(slide.shapes[SHAPE_TITLE], content.title)
    _set_text(slide.shapes[SHAPE_POSITION_LINE1], content.position_level)
    _set_text(slide.shapes[SHAPE_POSITION_LINE2], content.schwerpunkte)
    _set_text(slide.shapes[SHAPE_SUMMARY], content.summary)

    _set_bullets(slide.shapes[SHAPE_KOMPETENZEN], content.kompetenzen)
    _set_categorized_bullets(slide.shapes[SHAPE_RELEVANTE], content.relevante_erfahrungen)
    _set_bullets(slide.shapes[SHAPE_INTERNATIONAL], content.international_experience)
    _set_categorized_tools(slide.shapes[SHAPE_TOOLS], content.tool_categories)
    _set_bullets(slide.shapes[SHAPE_EDUCATION], content.education_certificates)

    if photo_path and Path(photo_path).exists():
        _replace_photo(slide, photo_path)

    prs.save(str(output_path))
    return output_path


def _set_text(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.size = Pt(BODY_FONT_PT)
        run.font.name = "+mn-lt"


def _set_bullets(shape, items: list[str]) -> None:
    tf = shape.text_frame
    tf.clear()
    for index, item in enumerate(items):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        if paragraph.runs:
            run = paragraph.runs[0]
            run.font.size = Pt(BODY_FONT_PT)
            run.font.name = "+mn-lt"
    _autofit_text_frame(tf)


def _set_categorized_bullets(shape, items) -> None:
    tf = shape.text_frame
    tf.clear()
    for index, item in enumerate(items):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.level = 0

        category_run = paragraph.add_run()
        category_run.text = f"{item.category}: "
        category_run.font.bold = True
        category_run.font.size = Pt(BODY_FONT_PT)
        category_run.font.name = "+mn-lt"

        detail_run = paragraph.add_run()
        detail_run.text = item.details
        detail_run.font.bold = False
        detail_run.font.size = Pt(BODY_FONT_PT)
        detail_run.font.name = "+mn-lt"
    _autofit_text_frame(tf)


def _set_categorized_tools(shape, categories) -> None:
    tf = shape.text_frame
    tf.clear()
    for index, category in enumerate(categories):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT

        label_run = paragraph.add_run()
        label_run.text = f"{category.category}: "
        label_run.font.bold = True
        label_run.font.size = Pt(BODY_FONT_PT)
        label_run.font.name = "+mn-lt"

        tools_run = paragraph.add_run()
        tools_run.text = ", ".join(category.tools)
        tools_run.font.bold = False
        tools_run.font.size = Pt(BODY_FONT_PT)
        tools_run.font.name = "+mn-lt"
    _autofit_text_frame(tf)


def _autofit_text_frame(text_frame) -> None:
    """Reduce font size if too many paragraphs are present."""
    paragraph_count = len([p for p in text_frame.paragraphs if p.text.strip()])
    if paragraph_count <= 6:
        return
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(MIN_FONT_PT)


def _replace_photo(slide, photo_path: str | Path) -> None:
    photo_path = Path(photo_path)
    shape = slide.shapes[SHAPE_PHOTO]
    left, top, width, height = shape.left, shape.top, shape.width, shape.height

    shape_element = shape._element
    shape_element.getparent().remove(shape_element)

    slide.shapes.add_picture(str(photo_path), left, top, width=width, height=height)
