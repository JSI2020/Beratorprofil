"""Verify PPTX formatting is preserved after generation."""

from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

from src.generator.pptx_generator import generate_pptx
from src.transformer.template_profiles import build_profile_content
from tests.test_template_alignment import _telecom_cv

PROJECT_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = PROJECT_ROOT / "template" / "Beraterprofil_TEMPLATE.pptx"

KOMPETENZEN_POS = (533999, 2943147)
SCHWERPUNKTE_POS = (3727239, 1363311)
SUMMARY_POS = (3727239, 1590943)
POSITION_POS = (3727239, 1050858)
TOOLS_POS = (553881, 4846422)


def _shape_at(slide, left: int, top: int):
    for shape in slide.shapes:
        if abs(shape.left - left) < 50000 and abs(shape.top - top) < 50000:
            return shape
    raise AssertionError(f"No shape at {left},{top}")


def test_kompetenzen_and_schwerpunkte_visible(tmp_path):
    content = build_profile_content(_telecom_cv(), "Funknetzplanung")
    output = tmp_path / "test.pptx"
    generate_pptx(content, TEMPLATE, output)

    slide = Presentation(output).slides[0]
    kompetenzen = _shape_at(slide, *KOMPETENZEN_POS)
    schwerpunkte = _shape_at(slide, *SCHWERPUNKTE_POS)

    assert kompetenzen.text_frame.text.strip()
    assert len(kompetenzen.text_frame.paragraphs) >= 6
    assert schwerpunkte.text_frame.text.strip()
    assert "Funknetz" in schwerpunkte.text_frame.text or "RAN" in schwerpunkte.text_frame.text


def test_summary_aligns_with_schwerpunkte(tmp_path):
    content = build_profile_content(_telecom_cv(), "Funknetzplanung")
    output = tmp_path / "test.pptx"
    generate_pptx(content, TEMPLATE, output)

    slide = Presentation(output).slides[0]
    schwerpunkte = _shape_at(slide, *SCHWERPUNKTE_POS)
    summary = _shape_at(slide, *SUMMARY_POS)

    assert summary.left == schwerpunkte.left
    assert summary.width == schwerpunkte.width
    assert summary.text_frame.text.strip()


def test_ausbildung_divider_below_text(tmp_path):
    content = build_profile_content(_telecom_cv(), "Funknetzplanung")
    output = tmp_path / "test.pptx"
    generate_pptx(content, TEMPLATE, output)

    slide = Presentation(output).slides[0]
    international = _shape_at(slide, 7997235, 2906571)
    divider = None
    for shape in slide.shapes:
        if shape.shape_type == 9 and shape.left > 7_000_000:
            divider = shape
            break

    assert divider is not None
    gap = divider.top - (international.top + international.height)
    assert gap >= 40_000
    assert gap <= 200_000


def test_tool_lines_keep_bold_category_runs(tmp_path):
    content = build_profile_content(_telecom_cv(), "Funknetzplanung")
    output = tmp_path / "test.pptx"
    generate_pptx(content, TEMPLATE, output)

    slide = Presentation(output).slides[0]
    tools = _shape_at(slide, *TOOLS_POS)
    first_para = tools.text_frame.paragraphs[0]
    assert first_para.runs[0].font.bold is True
    assert first_para.runs[1].font.bold is not True
