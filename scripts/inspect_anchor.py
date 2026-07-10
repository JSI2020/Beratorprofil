import sys
from pptx import Presentation
from pptx.oxml.ns import qn

path = sys.argv[1]
prs = Presentation(path)
slide = prs.slapes[0] if False else prs.slides[0]

for idx in [2, 10, 12, 13]:
    s = slide.shapes[idx]
    print(f"\n=== Shape {idx}: {s.name} ===")
    print(f"  pos: {s.left},{s.top} size: {s.width}x{s.height}")
    if hasattr(s, "fill"):
        try:
            print(f"  fill type: {s.fill.type}")
        except Exception as e:
            print(f"  fill: {e}")
    tf = s.text_frame
    body = tf._txBody
    body_pr = body.find(qn("a:bodyPr"))
    if body_pr is not None:
        print(f"  bodyPr: anchor={body_pr.get('anchor')} wrap={body_pr.get('wrap')} lIns={body_pr.get('lIns')}")
        for child in body_pr:
            print(f"    child: {child.tag.split('}')[-1]}")
