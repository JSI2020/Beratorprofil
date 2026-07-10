"""Extract layout profile coordinates from a manually tuned PPTX file.

Usage:
  python scripts/sync_layout_profile.py output/Farhan_Beraterprofil_LLM.pptx
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation

# Map template anchor positions to profile keys (original ORBIT template).
ANCHOR_MAP: dict[str, tuple[int, int]] = {
    "position_labels": (2209140, 972000),
    "summary_label": (2209140, 1569297),
    "kompetenzen_label": (647403, 2650668),
    "position": (3727239, 970730),
    "schwerpunkte": (3727239, 1291395),
    "summary": (3727239, 1590943),
    "kompetenzen": (533999, 2943147),
    "relevante": (4062770, 2906571),
    "international": (7997235, 2906571),
    "education": (8022387, 4720513),
    "tools": (553881, 4846422),
    "tools_dup": (438105, 4788720),
}

TOLERANCE = 120000


def _match_shape(slide, anchor_left: int, anchor_top: int):
    best = None
    best_dist = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        dist = abs(shape.left - anchor_left) + abs(shape.top - anchor_top)
        if dist > TOLERANCE:
            continue
        if best_dist is None or dist < best_dist:
            best = shape
            best_dist = dist
    return best


def extract_profile(path: Path) -> dict[str, tuple[int, int, int, int]]:
    slide = Presentation(str(path)).slides[0]
    profile: dict[str, tuple[int, int, int, int]] = {}
    for key, (left, top) in ANCHOR_MAP.items():
        shape = _match_shape(slide, left, top)
        if shape is None:
            print(f"Warning: could not match {key}", file=sys.stderr)
            continue
        profile[key] = (shape.left, shape.top, shape.width, shape.height)
    return profile


def format_profile_python(profile: dict[str, tuple[int, int, int, int]]) -> str:
    lines = ["LAYOUT_PROFILE: dict[str, ShapeRect] = {"]
    for key, (left, top, width, height) in profile.items():
        lines.append(f'    "{key}": ShapeRect({left}, {top}, {width}, {height}),')
    lines.append("}")
    return "\n".join(lines)


def main() -> int:
    if len(sys.argv) < 2:
        print("Usage: python scripts/sync_layout_profile.py <tuned.pptx>")
        return 1

    path = Path(sys.argv[1])
    profile = extract_profile(path)
    print(format_profile_python(profile))

    for shape in Presentation(str(path)).slides[0].shapes:
        if shape.shape_type == 9 and shape.left > 7_000_000:
            print(f"\n# RIGHT_COLUMN_LINE_TOP_HINT = {shape.top}")
            break
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
