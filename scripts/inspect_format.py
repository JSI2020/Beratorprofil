"""Inspect PPTX formatting details."""
from __future__ import annotations

import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


def inspect_shapes(path: str) -> None:
    prs = Presentation(path)
    slide = prs.slides[0]
    print(f"\nFILE: {path}")
    for idx in [2, 7, 12, 13, 14, 20]:
        shape = slide.shapes[idx]
        print(f"\n=== SHAPE {idx}: {shape.name} ===")
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for pi, para in enumerate(tf.paragraphs[:12]):
            text = para.text.replace("\n", " ").strip()
            if not text:
                continue
            print(f"  P{pi} level={para.level}: {text[:90]!r}")
            for ri, run in enumerate(para.runs):
                color = run.font.color
                color_info = "none"
                if color is not None and color.type is not None:
                    color_info = str(color.type)
                    try:
                        color_info += f" theme={color.theme_color}"
                    except Exception:
                        pass
                    try:
                        color_info += f" rgb={color.rgb}"
                    except Exception:
                        pass
                print(
                    f"      run{ri}: bold={run.font.bold} size={run.font.size} "
                    f"color={color_info} | {run.text[:55]!r}"
                )
            p_pr = para._p.find(qn("a:pPr"))
            if p_pr is not None:
                bu = p_pr.find(qn("a:buChar"))
                bu_auto = p_pr.find(qn("a:buAutoNum"))
                bu_font = p_pr.find(qn("a:buFont"))
                char = bu.get("char") if bu is not None else None
                print(
                    f"      bullet: char={char!r} auto={bu_auto is not None} "
                    f"buFont={bu_font.get('typeface') if bu_font is not None else None}"
                )


def inspect_xml_bullets(path: str) -> None:
    with zipfile.ZipFile(path) as zf:
        xml = zf.read("ppt/slides/slide1.xml").decode("utf-8", errors="ignore")
    chars = re.findall(r'buChar char="([^"]+)"', xml)
    from collections import Counter

    print(f"\nXML bullet chars in {path}: {Counter(chars)}")
    for char in sorted(set(chars)):
        print(f"  {char!r} U+{ord(char):04X}")


if __name__ == "__main__":
    target = sys.argv[1] if len(sys.argv) > 1 else r"template\Beraterprofil_TEMPLATE.pptx"
    inspect_shapes(target)
    inspect_xml_bullets(target)
