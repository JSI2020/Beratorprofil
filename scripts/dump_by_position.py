import sys
from pptx import Presentation

path = sys.argv[1]
prs = Presentation(path)
slide = prs.slides[0]

targets = {
    "kompetenzen": (533999, 2943147),
    "schwerpunkte": (3727239, 1291395),
    "position": (3727239, 970730),
}

for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    for name, (left, top) in targets.items():
        if abs(shape.left - left) < 50000 and abs(shape.top - top) < 50000:
            text = shape.text_frame.text.strip().replace("\n", " | ")
            print(f"{name}: z-order via iteration, text={text[:120]}")
