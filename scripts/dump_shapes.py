import sys
from pptx import Presentation

path = sys.argv[1]
prs = Presentation(path)
slide = prs.slides[0]
shapes = {
    2: "Kompetenzen",
    3: "Relevante",
    4: "Ausbildung/Karriere",
    12: "Position",
    13: "Schwerpunkte",
    14: "Summary",
    15: "Abschluss",
    20: "Tools",
}
for idx, name in shapes.items():
    s = slide.shapes[idx]
    t = s.text_frame.text.strip()
    print(f"[{idx}] {name}: paragraphs={len(s.text_frame.paragraphs)} chars={len(t)}")
    if t:
        print(t[:300])
    else:
        print("  *** EMPTY ***")
    print()
