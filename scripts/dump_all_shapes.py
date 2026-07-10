import sys
from pptx import Presentation

path = sys.argv[1]
slide = Presentation(path).slides[0]
for i, s in enumerate(slide.shapes):
    if not s.has_text_frame:
        print(f"[{i}] {s.shape_type} (no text) left={s.left} top={s.top}")
        continue
    t = s.text_frame.text.strip().replace("\n", " | ")[:120]
    empty = "*** EMPTY ***" if not t else t
    print(f"[{i}] left={s.left} top={s.top} w={s.width} h={s.height}")
    print(f"     {empty}")
