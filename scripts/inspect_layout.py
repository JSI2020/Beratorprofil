"""Inspect paragraph margins and line positions in template/output."""
import sys
from pptx import Presentation
from pptx.util import Inches

path = sys.argv[1]
slide = Presentation(path).slides[0]

print("=== Header text shapes ===")
for key, left, top in [
    ("position", 3727239, 970730),
    ("schwerpunkte", 3727239, 1291395),
    ("summary", 3727239, 1590943),
]:
    for s in slide.shapes:
        if s.has_text_frame and abs(s.left - left) < 50000 and abs(s.top - top) < 50000:
            tf = s.text_frame
            p = tf.paragraphs[0]
            print(f"{key}: left={s.left} top={s.top} w={s.width} h={s.height}")
            print(f"  margin_l={tf.margin_left} margin_r={tf.margin_right}")
            print(f"  p.level={p.level} indent={p.paragraph_format.left_indent} first={p.paragraph_format.first_line_indent}")
            print(f"  text={tf.text[:80]!r}")
            break

print("\n=== Lines ===")
for i, s in enumerate(slide.shapes):
    if str(s.shape_type) == "LINE (9)" or s.shape_type == 9:
        print(f"line[{i}] left={s.left} top={s.top} w={s.width} h={s.height}")

print("\n=== International (Ausbildung) ===")
for s in slide.shapes:
    if s.has_text_frame and abs(s.left - 7997235) < 50000 and abs(s.top - 2906571) < 50000:
        print(f"intl: left={s.left} top={s.top} w={s.width} h={s.height} bottom={s.top + s.height}")
        print(f"  paragraphs={len(s.text_frame.paragraphs)}")
        print(f"  text lines: {s.text_frame.text.count(chr(10)) + 1}")
