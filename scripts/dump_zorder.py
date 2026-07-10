import sys
from pptx import Presentation
from pptx.util import Inches

path = sys.argv[1]
prs = Presentation(path)
slide = prs.slides[0]

print(f"Total shapes: {len(slide.shapes)}\n")
for i, shape in enumerate(slide.shapes):
  if not shape.has_text_frame:
    kind = "no-text"
    preview = ""
  else:
    kind = "text"
    preview = shape.text_frame.text.strip().replace("\n", " | ")[:60]
  left = round(shape.left / 914400, 2)
  top = round(shape.top / 914400, 2)
  w = round(shape.width / 914400, 2)
  h = round(shape.height / 914400, 2)
  print(f"z={i:2d} {shape.name[:30]:30s} {left:5.2f},{top:5.2f} {w:5.2f}x{h:5.2f}  {preview}")
